"""Inspect specific challenging slides to verify quality."""
from pptx import Presentation
from pptx.util import Emu


def emu_to_in(v):
    return v / 914400 if v else 0


def show_slide(path, idx, label):
    print(f"\n--- {label} slide {idx} ---")
    prs = Presentation(path)
    s = prs.slides[idx - 1]
    print(f"  {len(s.shapes)} shapes:")
    for i, sh in enumerate(s.shapes):
        x = emu_to_in(sh.left); y = emu_to_in(sh.top)
        w = emu_to_in(sh.width); h = emu_to_in(sh.height)
        kind = ""
        if sh.has_text_frame:
            t = sh.text_frame.text.strip().replace('\n', ' | ')[:80]
            kind = f"tf={t!r}"
        elif sh.has_table:
            kind = f"table {len(sh.table.rows)}x{len(sh.table.columns)}"
        else:
            kind = "shape"
        print(f"  [{i:2d}] {sh.name:25s} pos=({x:.2f},{y:.2f}) size=({w:.2f}x{h:.2f}) {kind}")


# Day 1: cover, agenda, content slides
show_slide('day1-slides.pptx', 1, 'DAY 1 cover')
show_slide('day1-slides.pptx', 2, 'DAY 1 agenda (two-column)')
show_slide('day1-slides.pptx', 7, 'DAY 1 elevator pitch formula')
show_slide('day1-slides.pptx', 8, 'DAY 1 elevator pitch examples (cards-grid)')
show_slide('day1-slides.pptx', 13, 'DAY 1 C.R.E.A.T.E. prompt framework')
show_slide('day1-slides.pptx', 20, 'DAY 1 safety matrix table')
show_slide('day1-slides.pptx', 23, 'DAY 1 convergent thinking (two-column + prompt)')
show_slide('day1-slides.pptx', 24, 'DAY 1 SCAMPER intro (7 cards-grid)')
show_slide('day1-slides.pptx', 36, 'DAY 1 VPC alignment (two-column)')
show_slide('day1-slides.pptx', 43, 'DAY 1 wrap-up')

# Day 2
show_slide('day2-slides.pptx', 1, 'DAY 2 cover')
show_slide('day2-slides.pptx', 9, 'DAY 2 ideogram (highlight + prompt)')
show_slide('day2-slides.pptx', 25, 'DAY 2 custom domain (cards-grid)')
show_slide('day2-slides.pptx', 42, 'DAY 2 wrap-up')