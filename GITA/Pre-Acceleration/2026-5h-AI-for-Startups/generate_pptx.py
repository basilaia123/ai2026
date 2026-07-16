"""Generate Day 1 and Day 2 PPTX files from parsed JSON.

Uses python-pptx. Theme matches the source HTML design:
- Primary: amber (#f59e0b)
- Dark accent: amber-600 (#d97706)
- Status: red / amber / green / blue
"""
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- Theme ----------
THEME = {
    "primary": RGBColor(0x4F, 0x46, 0xE5),       # indigo-600
    "primary_dark": RGBColor(0x37, 0x30, 0xA3),  # indigo-800
    "primary_light": RGBColor(0xEE, 0xF2, 0xFF), # indigo-50
    "accent": RGBColor(0xF4, 0x3F, 0x5E),        # rose-500
    "success": RGBColor(0x05, 0x96, 0x69),       # emerald-600
    "info": RGBColor(0x3B, 0x82, 0xF6),          # blue-500
    "warning": RGBColor(0x4F, 0x46, 0xE5),       # alias to primary
    "danger": RGBColor(0xF4, 0x3F, 0x5E),
    "text": RGBColor(0x0F, 0x17, 0x2A),          # slate-900
    "text_secondary": RGBColor(0x47, 0x55, 0x69),# slate-600
    "text_muted": RGBColor(0x94, 0xA3, 0xB8),    # slate-400
    "bg": RGBColor(0xFF, 0xFF, 0xFF),
    "bg_light": RGBColor(0xF8, 0xFA, 0xFC),      # slate-50
    "border": RGBColor(0xE2, 0xE8, 0xF0),        # slate-200
    "highlight": RGBColor(0xEE, 0xF2, 0xFF),     # indigo-50
    "warning_bg": RGBColor(0xFF, 0xF7, 0xED),    # orange-50
    "danger_bg": RGBColor(0xFF, 0xF1, 0xF2),     # rose-50
    "success_bg": RGBColor(0xEC, 0xFD, 0xF5),    # emerald-50
    "info_bg": RGBColor(0xEF, 0xF6, 0xFF),       # blue-50
}

FONT_TITLE = "Sylfaen"   # ships with Windows; renders Georgian well
FONT_BODY = "Sylfaen"
FONT_MONO = "Consolas"


# ---------- Layout ----------
SLIDE_W = 10.0
SLIDE_H = 5.625


def set_run(run, text, *, size=None, bold=None, italic=None, color=None, font=None):
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    if font is not None:
        run.font.name = font
        # also set east-asian font for safety
        rPr = run._r.get_or_add_rPr()
        for tag in ("ea", "cs", "latin"):
            existing = rPr.find(qn(f"a:{tag}"))
            if existing is not None:
                rPr.remove(existing)
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
            el.set("typeface", font)


def add_rect(slide, x, y, w, h, *, fill, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, *, fill, line=None, line_w=0, radius=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    # Adjust corner radius
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def set_solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def add_textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, word_wrap=True, margin=0.06):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    return tb, tf


def add_paragraph(tf, text, *, size=14, bold=False, italic=False, color=None, font=None, align=None, space_after=4):
    if tf.paragraphs and tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    set_run(p.add_run(), text, size=size, bold=bold, italic=italic,
            color=color or THEME["text"], font=font or FONT_BODY)
    return p


def add_runs_paragraph(tf, runs, *, align=None, space_after=4):
    if tf.paragraphs and tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    for r in runs:
        set_run(p.add_run(), **r)
    return p


# ---------- Header / Footer chrome ----------
def draw_header(slide, slide_num, total, title_text):
    # Top amber bar
    add_rect(slide, 0, 0, SLIDE_W, 0.12, fill=THEME["primary"])

    # Title text
    if title_text:
        title_tb, tf = add_textbox(slide, 0.5, 0.28, 8.5, 0.85, anchor=MSO_ANCHOR.TOP)
        # Strip emojis from title for cleaner look (keep some though)
        clean_title = title_text
        # Allow emojis but don't let them break sizing
        run = tf.paragraphs[0].add_run()
        set_run(run, clean_title, size=26, bold=True,
                color=THEME["primary_dark"], font=FONT_TITLE)
        # Underline rule under title
        add_rect(slide, 0.5, 1.15, 9.0, 0.025, fill=THEME["border"])

    # Page number badge top-right
    if slide_num and total:
        badge_w = 0.85
        badge_x = SLIDE_W - badge_w - 0.25
        badge = add_rounded(slide, badge_x, 0.22, badge_w, 0.35,
                            fill=THEME["primary"], line=None, radius=0.5)
        tf = badge.text_frame
        tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), f"{slide_num} / {total}",
                size=10, bold=True, color=THEME["bg"], font=FONT_BODY)


def draw_footer(slide, *, footer_text="GITA Pre-Acceleration 2026 · AI for Startups"):
    # bottom thin line + footer text
    add_rect(slide, 0.5, SLIDE_H - 0.32, 9.0, 0.015, fill=THEME["border"])
    tb, tf = add_textbox(slide, 0.5, SLIDE_H - 0.28, 9.0, 0.22, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraph(tf, footer_text, size=8, color=THEME["text_muted"], font=FONT_BODY)


# ---------- Slide renderers ----------
def make_blank_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def add_cover_slide(prs, slide_data, total):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Background fill
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=THEME["bg"])

    # Top amber bar
    add_rect(slide, 0, 0, SLIDE_W, 0.55, fill=THEME["primary"])

    # Decorative side strip
    add_rect(slide, 0, 0, 0.18, SLIDE_H, fill=THEME["primary_dark"])

    # Title block
    title = slide_data.get("title") or ""
    tb, tf = add_textbox(slide, 0.7, 1.4, 8.6, 1.4, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraph(tf, title, size=36, bold=True,
                  color=THEME["primary_dark"], font=FONT_TITLE, space_after=4)

    # Subtitle / kicker from first paragraph if available
    blocks = slide_data.get("blocks", [])
    kicker = ""
    bullets = []
    extra_blocks = []
    seen_first_p = False
    if blocks:
        for b in blocks:
            if b.get("type") == "p" and not seen_first_p:
                kicker = b["text"]
                seen_first_p = True
            elif b.get("type") == "list":
                bullets = b["items"][:4]
            elif b.get("type") == "box" and "key-points" in (b.get("class") or ""):
                # Extract h4 + paragraphs from key-points box
                kp_h4 = None
                kp_paras = []
                for sb in b.get("blocks", []):
                    if sb.get("type") == "h4":
                        kp_h4 = sb.get("text")
                    elif sb.get("type") == "p":
                        kp_paras.append(sb.get("text"))
                if kp_h4:
                    extra_blocks.append(("h4", kp_h4))
                for kp in kp_paras:
                    extra_blocks.append(("p", kp))

    cur_y = 2.85
    if kicker:
        tb, tf = add_textbox(slide, 0.7, cur_y, 8.6, 0.6, anchor=MSO_ANCHOR.TOP)
        add_paragraph(tf, kicker, size=16, color=THEME["text_secondary"], font=FONT_BODY)
        cur_y += 0.7

    if bullets:
        tb, tf = add_textbox(slide, 0.7, cur_y, 8.6, 1.4, anchor=MSO_ANCHOR.TOP)
        for it in bullets:
            add_paragraph(tf, "• " + it, size=13, color=THEME["text"], font=FONT_BODY, space_after=2)
        cur_y += 1.5

    # Render key-points content
    for kind, text in extra_blocks:
        if cur_y > SLIDE_H - 0.8:
            break
        if kind == "h4":
            tb, tf = add_textbox(slide, 0.7, cur_y, 8.6, 0.32, anchor=MSO_ANCHOR.TOP)
            add_paragraph(tf, text, size=13, bold=True,
                          color=THEME["primary_dark"], font=FONT_BODY)
            cur_y += 0.34
        elif kind == "p":
            tb, tf = add_textbox(slide, 0.7, cur_y, 8.6, 0.30, anchor=MSO_ANCHOR.TOP)
            add_paragraph(tf, text, size=12, color=THEME["text"], font=FONT_BODY)
            cur_y += 0.32

    # Bottom branding
    tb, tf = add_textbox(slide, 0.7, SLIDE_H - 0.55, 8.6, 0.3, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraph(tf, "GITA Pre-Acceleration 2026 · AI for Startups · გიორგი ბასილაია",
                  size=10, color=THEME["text_muted"], font=FONT_BODY, align=PP_ALIGN.CENTER)


def add_content_slide(prs, slide_data, total, slide_index):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=THEME["bg"])

    # Left amber strip
    add_rect(slide, 0, 0, 0.12, SLIDE_H, fill=THEME["primary"])

    # Title + page number
    draw_header(slide, slide_index, total, slide_data.get("title") or "")

    # Content area
    content_y = 1.3
    content_h = SLIDE_H - content_y - 0.32  # leave footer room (footer is ~0.32 tall)
    blocks = slide_data.get("blocks", [])
    render_blocks(slide, blocks, 0.5, content_y, 9.0, content_h)

    draw_footer(slide)


def render_blocks(slide, blocks, x, y, w, h):
    """Render a list of content blocks within the (x,y,w,h) area.
    Uses a simple vertical flow; each block gets a max height proportional to its content.
    Strict bounds: no overflow allowed beyond content area, but small prompt-boxes always render.
    """
    cur_y = y
    remaining_h = h

    for block in blocks:
        if cur_y >= y + h - 0.02:
            break  # stop before exceeding area

        btype = block.get("type")
        cls = block.get("class", "")
        # Prompt-boxes always render even when tight (they're important code templates)
        is_prompt_box = btype == "box" and "prompt" in cls

        if remaining_h < 0.10 and not is_prompt_box:
            break

        min_h = _min_block_height(block, w)
        pref_h = estimate_block_height(block, w)
        block_h = min(pref_h, remaining_h)

        # For prompt-boxes, force rendering even at very small heights
        if is_prompt_box and block_h < min_h:
            block_h = max(min_h, min(0.50, remaining_h))

        if btype == "p":
            render_paragraph(slide, block, x, cur_y, w, block_h)
        elif btype == "h3":
            render_h3(slide, block, x, cur_y, w, block_h)
        elif btype == "h4":
            render_h4(slide, block, x, cur_y, w, block_h)
        elif btype == "list":
            render_list(slide, block, x, cur_y, w, block_h)
        elif btype == "table":
            render_table(slide, block, x, cur_y, w, block_h)
        elif btype == "code":
            render_code(slide, block, x, cur_y, w, block_h)
        elif btype == "box":
            render_box(slide, block, x, cur_y, w, block_h)
        else:
            tb, tf = add_textbox(slide, x, cur_y, w, block_h)
            add_paragraph(tf, str(block.get("text", "")), size=12, color=THEME["text"], font=FONT_BODY)

        cur_y += block_h + 0.08
        remaining_h -= block_h + 0.08


def _min_block_height(block, w):
    """Minimum height required to render the block at all (one line)."""
    btype = block.get("type")
    if btype == "p":
        return 0.32
    if btype == "h3":
        return 0.40
    if btype == "h4":
        return 0.32
    if btype == "list":
        return 0.32
    if btype == "table":
        return 0.45
    if btype == "code":
        return 0.42
    if btype == "box":
        cls = block.get("class", "")
        sub = block.get("blocks") or []
        # Compact prompt-box has its own compact minimum
        if "prompt" in cls and len(sub) >= 2:
            return 0.30  # very compact: label + tiny code area
        if "prompt-label" in cls and len(sub) == 1:
            return 0.30
        if sub:
            return sum(_min_block_height(b, w) for b in sub) + 0.1
        return 0.35
    return 0.3


def estimate_block_height(block, w):
    btype = block.get("type")
    if btype == "h3":
        return 0.45
    if btype == "h4":
        return 0.35
    if btype == "p":
        text = block.get("text", "")
        chars = len(text)
        # roughly 80 chars per line at 13pt over ~9 inches
        lines = max(1, (chars + 79) // 80)
        return min(0.26 + 0.26 * lines, 1.5)
    if btype == "list":
        n = len(block.get("items", []))
        if n <= 0:
            return 0.2
        # Estimate chars per line based on width: ~85 chars per inch at 13pt
        chars_per_line = max(20, int(w * 9))  # roughly 9 chars/inch usable
        total = 0
        for it in block["items"]:
            chars = len(it)
            lines = max(1, (chars + chars_per_line - 1) // chars_per_line)
            total += 0.16 + 0.30 * lines
        return min(total + 0.2, 4.5)
    if btype == "table":
        rows = block.get("rows", [])
        n = len(rows)
        return min(0.45 + 0.35 * n, 3.5)
    if btype == "code":
        text = block.get("text", "")
        # Estimate wrapped lines based on width
        char_width = 7.0
        chars_per_line = max(20, int((w - 0.32) * char_width))
        total_lines = 0
        for raw_line in text.split("\n"):
            if not raw_line:
                total_lines += 1
                continue
            n = max(1, (len(raw_line) + chars_per_line - 1) // chars_per_line)
            total_lines += n
        return min(0.30 + 0.22 * total_lines + 0.10, 3.5)
    if btype == "box":
        cls = block.get("class", "")
        sub = block.get("blocks") or []
        if sub:
            # For multi-column layouts, the available width per child is smaller,
            # so list items wrap to more lines. Use narrower width in estimate.
            if "two-column" in cls:
                child_w = (w - 0.4) / 2
                # height = max child (both columns same height)
                return max(estimate_block_height(b, child_w) for b in sub) + 0.3
            elif "cards-grid" in cls:
                if len(sub) <= 2:
                    child_w = (w - 0.4) / 2
                elif len(sub) <= 4:
                    child_w = (w - 0.4) / 2
                elif len(sub) <= 6:
                    child_w = (w - 0.4) / 3
                else:
                    child_w = (w - 0.4) / 4
                # height = sum of rows
                n_cols = 2 if len(sub) <= 4 else (3 if len(sub) <= 6 else 4)
                n_rows = (len(sub) + n_cols - 1) // n_cols
                row_heights = []
                for r in range(n_rows):
                    rs = sub[r * n_cols:(r + 1) * n_cols]
                    row_heights.append(max(estimate_block_height(b, child_w) for b in rs))
                return sum(row_heights) + 0.15 * (n_rows - 1) + 0.3
            else:
                child_w = w - 0.4
                # Tight padding for compact label boxes (prompt-label)
                if "prompt-label" in cls and len(sub) == 1 and sub[0].get("type") == "code":
                    return 0.32
                # Tight compact for prompt-box (combined header + code)
                if "prompt" in cls and len(sub) >= 2:
                    label_h = 0.32
                    code_h = 0.20 + 0.20 * 3 + 0.10  # estimate 3-line code
                    return label_h + 0.06 + code_h + 0.10
                # Use minimal padding when box wraps a single inline element
                pad = 0.1 if len(sub) == 1 else 0.3
                return sum(estimate_block_height(b, child_w) for b in sub) + pad
        text = block.get("text", "")
        chars = len(text)
        lines = max(1, (chars + 90) // 90)
        return min(0.35 + 0.30 * lines, 2.5)
    return 0.4


def render_paragraph(slide, block, x, y, w, h):
    tb, tf = add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP)
    # Use 13pt for dense slides to save vertical space
    add_paragraph(tf, block["text"], size=13, color=THEME["text"], font=FONT_BODY)


def render_h3(slide, block, x, y, w, h):
    tb, tf = add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP)
    add_paragraph(tf, block["text"], size=18, bold=True, color=THEME["primary_dark"], font=FONT_TITLE)


def render_h4(slide, block, x, y, w, h):
    tb, tf = add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP)
    add_paragraph(tf, block["text"], size=14, bold=True, color=THEME["text"], font=FONT_BODY)


def render_list(slide, block, x, y, w, h):
    ordered = block.get("ordered", False)
    items = block.get("items", [])
    tb, tf = add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP)
    for i, item in enumerate(items):
        marker = f"{i+1}." if ordered else "▸"
        # First paragraph already exists
        if i == 0 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(3)
            set_run(p.add_run(), f"{marker}  ", size=13, bold=True,
                    color=THEME["primary_dark"], font=FONT_BODY)
            set_run(p.add_run(), item, size=13, color=THEME["text"], font=FONT_BODY)
        else:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(3)
            set_run(p.add_run(), f"{marker}  ", size=13, bold=True,
                    color=THEME["primary_dark"], font=FONT_BODY)
            set_run(p.add_run(), item, size=13, color=THEME["text"], font=FONT_BODY)


def render_table(slide, block, x, y, w, h):
    rows = block.get("rows", [])
    if not rows:
        return
    n_cols = max(len(r) for r in rows)
    n_rows = len(rows)
    # Compute column widths
    col_w = w / n_cols
    # Table dimensions
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                       Inches(w), Inches(min(h, 0.4 + 0.4 * n_rows)))
    tbl = tbl_shape.table
    for ci in range(n_cols):
        tbl.columns[ci].width = Inches(col_w)
    for ri in range(n_rows):
        row_h = min(0.5, h / n_rows) if n_rows > 0 else 0.5
        tbl.rows[ri].height = Inches(row_h)

    # Header row
    for ci, cell_text in enumerate(rows[0]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME["primary_dark"]
        tf = cell.text_frame
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        set_run(p.add_run(), cell_text, size=11, bold=True,
                color=THEME["bg"], font=FONT_BODY)

    # Data rows
    for ri in range(1, n_rows):
        # Decide row fill based on content (safety rating detection)
        row_text = " ".join(rows[ri]).upper()
        if "GO" in row_text or "🟢" in row_text or "D1FAE5" in row_text or "უსაფრთხო" in row_text.lower():
            fill = THEME["success_bg"]
        elif "STOP" in row_text or "🔴" in row_text or "წითელი" in row_text or "გამორიცხული" in row_text.lower():
            fill = THEME["danger_bg"]
        elif "CAUTION" in row_text or "🟡" in row_text or "ფრთხილად" in row_text.lower():
            fill = THEME["warning_bg"]
        else:
            fill = THEME["bg_light"] if ri % 2 == 0 else THEME["bg"]

        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            tf = cell.text_frame
            tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            tf.clear()
            text = rows[ri][ci] if ci < len(rows[ri]) else ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            # detect leading emoji or status
            set_run(p.add_run(), text, size=10, color=THEME["text"], font=FONT_BODY)


def render_code(slide, block, x, y, w, h):
    # Estimate wrapped lines based on width (Consolas 11pt: ~7 chars/inch)
    text = block["text"]
    char_width = 7.0  # approximate chars per inch for Consolas 11pt
    chars_per_line = max(20, int((w - 0.32) * char_width))
    # Wrap text into lines
    wrapped_lines = []
    for raw_line in text.split("\n"):
        if not raw_line:
            wrapped_lines.append("")
            continue
        # simple word-wrap
        words = raw_line.split(" ")
        cur = ""
        for w_word in words:
            if len(cur) + len(w_word) + 1 <= chars_per_line:
                cur = (cur + " " + w_word).strip()
            else:
                if cur:
                    wrapped_lines.append(cur)
                cur = w_word
        if cur:
            wrapped_lines.append(cur)
    n_lines = max(1, len(wrapped_lines))
    line_h = 0.22
    needed_h = 0.20 + line_h * n_lines + 0.10  # padding top + lines + padding bottom
    used_h = min(h, needed_h)

    # Outer box
    add_rounded(slide, x, y, w, used_h,
                fill=THEME["bg_light"], line=THEME["border"], line_w=0.75, radius=0.04)
    # Inner code area
    inner_y = y + 0.05
    inner_h = used_h - 0.10
    add_rect(slide, x + 0.06, inner_y, w - 0.12, inner_h,
             fill=THEME["bg"], line=THEME["border"], line_w=0.5)
    tb, tf = add_textbox(slide, x + 0.12, inner_y + 0.02, w - 0.24, inner_h - 0.04, anchor=MSO_ANCHOR.TOP)
    tf.word_wrap = True
    for i, line in enumerate(wrapped_lines):
        if i == 0 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(0)
        set_run(p.add_run(), line if line else " ", size=10,
                color=THEME["text"], font=FONT_MONO)


def _box_palette(cls):
    """Return (fill, line, accent) for a given box class string."""
    if "highlight" in cls:
        return THEME["primary_light"], THEME["primary_dark"], THEME["primary_dark"]
    if "warning" in cls:
        return THEME["warning_bg"], THEME["warning"], THEME["warning"]
    if "danger" in cls:
        return THEME["danger_bg"], THEME["danger"], THEME["danger"]
    if "success" in cls:
        return THEME["success_bg"], THEME["success"], THEME["success"]
    if "info" in cls:
        return THEME["info_bg"], THEME["info"], THEME["info"]
    if "prompt" in cls:
        return THEME["bg_light"], THEME["primary_dark"], THEME["primary_dark"]
    return THEME["bg_light"], THEME["border"], THEME["primary_dark"]


def render_box(slide, block, x, y, w, h):
    cls = block.get("class", "")
    fill, line, accent = _box_palette(cls)

    # Layout pattern dispatch
    if "two-column" in cls and "blocks" in block and len(block["blocks"]) >= 2:
        return _render_two_column(slide, block, x, y, w, h, fill, line, accent)
    if "cards-grid" in cls and "blocks" in block and len(block["blocks"]) >= 2:
        return _render_cards_grid(slide, block, x, y, w, h, fill, line, accent)

    # Compact prompt-box: combine label + code into a single block
    if "prompt" in cls and "blocks" in block:
        label_text = None
        code_text = None
        for sb in block["blocks"]:
            if sb.get("type") == "box" and "prompt-label" in (sb.get("class") or ""):
                # extract label
                for ssb in sb.get("blocks", []):
                    if ssb.get("type") == "code":
                        label_text = ssb.get("text", "")
                        break
            elif sb.get("type") == "code":
                code_text = sb.get("text", "")
        if label_text is not None or code_text is not None:
            _render_prompt_box_compact(slide, label_text or "", code_text or "", x, y, w, h)
            return

    # Compact rendering for label-like boxes (single code block, no decorative padding)
    if "prompt-label" in cls and "blocks" in block and len(block["blocks"]) == 1:
        sub = block["blocks"][0]
        if sub.get("type") == "code":
            # Just render the label code tightly
            tb, tf = add_textbox(slide, x, y, w, 0.42, anchor=MSO_ANCHOR.MIDDLE)
            tf.margin_left = Inches(0.10)
            tf.margin_right = Inches(0.10)
            add_paragraph(tf, sub.get("text", ""), size=11, bold=True,
                          color=THEME["primary_dark"], font=FONT_MONO)
            return

    # Default: single stacked box
    body_h = h
    add_rect(slide, x, y, 0.06, body_h, fill=accent)
    add_rounded(slide, x + 0.06, y, w - 0.06, body_h,
                fill=fill, line=line, line_w=0.75, radius=0.04)

    inner_x = x + 0.18
    inner_y = y + 0.08
    inner_w = w - 0.30
    inner_h = body_h - 0.16

    if "blocks" in block and block["blocks"]:
        render_blocks(slide, block["blocks"], inner_x, inner_y, inner_w, inner_h)
    elif "text" in block and block["text"]:
        tb, tf = add_textbox(slide, x + 0.18, inner_y, inner_w, inner_h, anchor=MSO_ANCHOR.TOP)
        add_paragraph(tf, block["text"], size=13, color=THEME["text"], font=FONT_BODY)


def _render_prompt_box_compact(slide, label_text, code_text, x, y, w, h):
    """Compact prompt box: single accent + inline label + code body. No nested boxes.
    Auto-fits to allocated height; minimum 0.30 inches.
    """
    # Clamp to minimum height
    used_h = max(0.30, min(h, 0.85))
    # Outer accent strip
    add_rect(slide, x, y, 0.06, used_h, fill=THEME["primary_dark"])
    # Outer rounded background
    add_rounded(slide, x + 0.06, y, w - 0.06, used_h,
                fill=THEME["bg_light"], line=THEME["border"], line_w=0.75, radius=0.04)

    # Inline label as a single line at the top (no separate bar)
    label_h = 0.24 if label_text else 0.0
    if label_text:
        tb, tf = add_textbox(slide, x + 0.12, y + 0.04, w - 0.24, label_h, anchor=MSO_ANCHOR.TOP)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        add_paragraph(tf, label_text, size=9, bold=True,
                      color=THEME["primary_dark"], font=FONT_MONO)

    # Code body fills remaining space
    code_y = y + 0.04 + label_h + 0.02
    code_h = used_h - 0.04 - label_h - 0.06
    if code_h < 0.10:
        code_h = 0.10  # minimum
    add_rect(slide, x + 0.12, code_y, w - 0.24, code_h,
             fill=THEME["bg"], line=THEME["border"], line_w=0.5)
    tb, tf = add_textbox(slide, x + 0.18, code_y + 0.03, w - 0.36, code_h - 0.06, anchor=MSO_ANCHOR.TOP)
    tf.word_wrap = True
    # Wrap code text
    char_width = 7.0
    chars_per_line = max(20, int((w - 0.36) * char_width))
    wrapped = []
    for raw in code_text.split("\n"):
        if not raw:
            wrapped.append("")
            continue
        words = raw.split(" ")
        cur = ""
        for wd in words:
            if len(cur) + len(wd) + 1 <= chars_per_line:
                cur = (cur + " " + wd).strip()
            else:
                if cur:
                    wrapped.append(cur)
                cur = wd
        if cur:
            wrapped.append(cur)
    line_h = 0.18
    avail_lines = max(1, int((code_h - 0.06) / line_h))
    if len(wrapped) > avail_lines:
        line_h = max(0.12, (code_h - 0.06) / max(1, len(wrapped)))
    for i, line in enumerate(wrapped):
        if i == 0 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(0)
        size = max(7, int(line_h * 72))
        set_run(p.add_run(), line if line else " ", size=size,
                color=THEME["text"], font=FONT_MONO)


def _render_two_column(slide, block, x, y, w, h, fill, line, accent):
    """Render side-by-side columns. Each column gets a sub-box."""
    subs = block["blocks"]
    n = len(subs)
    gap = 0.15
    col_w = (w - gap * (n - 1)) / n
    for i, sb in enumerate(subs):
        cx = x + i * (col_w + gap)
        sub_cls = sb.get("class", "")
        sub_fill, sub_line, sub_accent = _box_palette(sub_cls)
        # subtle outer card for the column
        add_rect(slide, cx, y, 0.06, h, fill=sub_accent)
        add_rounded(slide, cx + 0.06, y, col_w - 0.06, h,
                    fill=sub_fill, line=sub_line, line_w=0.5, radius=0.04)
        ix = cx + 0.18
        iy = y + 0.08
        iw = col_w - 0.30
        ih = h - 0.16
        if "blocks" in sb and sb["blocks"]:
            render_blocks(slide, sb["blocks"], ix, iy, iw, ih)
        elif "text" in sb and sb["text"]:
            tb, tf = add_textbox(slide, ix, iy, iw, ih, anchor=MSO_ANCHOR.TOP)
            add_paragraph(tf, sb["text"], size=12, color=THEME["text"], font=FONT_BODY)


def _render_cards_grid(slide, block, x, y, w, h, fill, line, accent):
    """Render N cards in a grid. 2 cols by default; 3 if there are 4+ cards."""
    subs = block["blocks"]
    n = len(subs)
    if n <= 2:
        n_cols = n
    elif n <= 4:
        n_cols = 2
    elif n <= 6:
        n_cols = 3
    else:
        n_cols = 4
    n_rows = (n + n_cols - 1) // n_cols
    gap_x = 0.15
    gap_y = 0.15
    cell_w = (w - gap_x * (n_cols - 1)) / n_cols
    cell_h = (h - gap_y * (n_rows - 1)) / n_rows
    for i, sb in enumerate(subs):
        r = i // n_cols
        c = i % n_cols
        cx = x + c * (cell_w + gap_x)
        cy = y + r * (cell_h + gap_y)
        sub_cls = sb.get("class", "")
        sub_fill, sub_line, sub_accent = _box_palette(sub_cls)
        add_rect(slide, cx, cy, 0.05, cell_h, fill=sub_accent)
        add_rounded(slide, cx + 0.05, cy, cell_w - 0.05, cell_h,
                    fill=sub_fill, line=sub_line, line_w=0.5, radius=0.04)
        ix = cx + 0.14
        iy = cy + 0.06
        iw = cell_w - 0.22
        ih = cell_h - 0.12
        if "blocks" in sb and sb["blocks"]:
            render_blocks(slide, sb["blocks"], ix, iy, iw, ih)
        elif "text" in sb and sb["text"]:
            tb, tf = add_textbox(slide, ix, iy, iw, ih, anchor=MSO_ANCHOR.TOP)
            add_paragraph(tf, sb["text"], size=11, color=THEME["text"], font=FONT_BODY)


def add_agenda_slide(prs, slide_data, total, slide_index):
    """Render agenda-style slide with numbered list, more visual than generic content."""
    add_content_slide(prs, slide_data, total, slide_index)


def add_summary_slide(prs, slide_data, total, slide_index):
    """Summary / closing slide."""
    add_content_slide(prs, slide_data, total, slide_index)


# ---------- Main pipeline ----------
def generate_pptx(json_path: Path, output_path: Path, deck_label: str):
    slides = json.loads(json_path.read_text(encoding="utf-8"))
    total = len(slides)
    print(f"\n=== {deck_label}: {total} slides ===")
    prs = make_blank_presentation()

    for s in slides:
        idx = s["index"]
        if idx == 1:
            add_cover_slide(prs, s, total)
        else:
            add_content_slide(prs, s, total, idx)

    prs.save(str(output_path))
    print(f"  saved: {output_path}  ({output_path.stat().st_size / 1024:.1f} KB)")


def main():
    workspace = Path(r"C:\Users\GBASILAIA\claude\ai2026\GITA\Pre-Acceleration\2026-5h-AI-for-Startups")
    generate_pptx(
        workspace / "day1-slides.json",
        workspace / "day1-slides.pptx",
        "Day 1 — AI Foundation & Validation",
    )
    generate_pptx(
        workspace / "day2-slides.json",
        workspace / "day2-slides.pptx",
        "Day 2 — Product & Launch",
    )


if __name__ == "__main__":
    main()