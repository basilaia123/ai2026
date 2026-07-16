"""Comprehensive QA: verify both PPTX files have correct structure and content."""
import json
from collections import Counter
from pptx import Presentation
from pptx.util import Emu


def emu_to_in(v):
    return v / 914400 if v else 0


def audit(path, json_path, label):
    print(f"\n{'='*70}")
    print(f"  AUDIT: {label}")
    print(f"  PPTX: {path}")
    print(f"{'='*70}")
    prs = Presentation(path)
    src = json.load(open(json_path, encoding='utf-8'))

    print(f"  Slides in PPTX: {len(prs.slides)}")
    print(f"  Slides in JSON: {len(src)}")
    print(f"  Slide size: {emu_to_in(prs.slide_width):.2f} x {emu_to_in(prs.slide_height):.2f} inches")

    if len(prs.slides) != len(src):
        print(f"  ⚠ MISMATCH: PPTX and JSON slide counts differ")
        return False

    issues = []
    text_chars = 0
    tables = 0
    for i, (pslide, jslide) in enumerate(zip(prs.slides, src), 1):
        title = jslide.get('title', '')
        # Collect all text on slide
        slide_text = ""
        for shape in pslide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text_frame.text + " "
            elif shape.has_table:
                tables += 1
                for row in shape.table.rows:
                    for cell in row.cells:
                        slide_text += cell.text_frame.text + " "
        text_chars += len(slide_text)

        # Check 1: Slide should have non-empty text
        if len(slide_text.strip()) < 20:
            issues.append(f"  Slide {i}: very little text ({len(slide_text)} chars)")

        # Check 2: page number badge on non-cover slides
        if i > 1:
            has_pn = f"{i} /" in slide_text or f"{i}/" in slide_text
            if not has_pn:
                issues.append(f"  Slide {i}: missing page number badge")

        # Check 3: source title should appear somewhere on slide
        if title and i > 1:
            # Allow partial matches for truncated titles
            title_words = title.split()[:3]
            if title_words and not any(w in slide_text for w in title_words):
                issues.append(f"  Slide {i}: title not found in slide text — title={title[:50]!r}")

        # Check 4: shapes fit within slide bounds
        for shape in pslide.shapes:
            x = emu_to_in(shape.left)
            y = emu_to_in(shape.top)
            w = emu_to_in(shape.width)
            h = emu_to_in(shape.height)
            if x + w > 10.05:
                issues.append(f"  Slide {i}: shape overflows right edge (x+w={x+w:.2f})")
            if y + h > 5.65:
                issues.append(f"  Slide {i}: shape overflows bottom edge (y+h={y+h:.2f})")

    print(f"  Total text characters: {text_chars:,}")
    print(f"  Total tables: {tables}")
    if issues:
        print(f"  ⚠ {len(issues)} issues found:")
        for iss in issues[:20]:
            print(iss)
        if len(issues) > 20:
            print(f"    ... and {len(issues) - 20} more")
    else:
        print(f"  ✓ No issues found")
    return len(issues) == 0


ok1 = audit('day1-slides.pptx', 'day1-slides.json', 'Day 1 — AI Foundation & Validation')
ok2 = audit('day2-slides.pptx', 'day2-slides.json', 'Day 2 — Product & Launch')
print(f"\n{'='*70}")
print(f"  RESULT: Day 1 {'PASS' if ok1 else 'FAIL'}  |  Day 2 {'PASS' if ok2 else 'FAIL'}")
print(f"{'='*70}")