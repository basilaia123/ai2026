from pptx import Presentation
for path, idx, label in [
    ('day1-slides.pptx', 31, 'Day 1'),
    ('day1-slides.pptx', 37, 'Day 1'),
    ('day2-slides.pptx', 21, 'Day 2'),
]:
    prs = Presentation(path)
    s = prs.slides[idx - 1]
    print(f"\n--- {label} slide {idx} ({len(s.shapes)} shapes) ---")
    for i, sh in enumerate(s.shapes):
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip().replace(chr(10), ' | ')[:60]
        else:
            txt = ''
        x = sh.left/914400; y = sh.top/914400; w = sh.width/914400; h = sh.height/914400
        print(f'  [{i:2d}] pos=({x:.2f},{y:.2f}) size=({w:.2f}x{h:.2f}) {txt!r}')