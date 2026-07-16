import os
from pptx import Presentation
for path, label, total in [('day1-slides.pptx', 'DAY 1', 43), ('day2-slides.pptx', 'DAY 2', 42)]:
    prs = Presentation(path)
    print(f'\n=== {label} — {total} slides — {len(prs.slides)} slides confirmed ===')
    print(f'File size: {os.path.getsize(path)/1024:.1f} KB')
    print(f'Slide size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches (16:9)')
    for idx, tag in [(1, 'Cover'), (2, 'Agenda'), (3, 'Welcome showcase 1'), (13, 'C.R.E.A.T.E.'), (total, 'Wrap-up')]:
        s = prs.slides[idx - 1]
        title_text = '(no title)'
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title_text = sh.text_frame.text.strip()[:50]
                break
        print(f'  Slide {idx:2d} [{tag:14s}]: {title_text!r}')