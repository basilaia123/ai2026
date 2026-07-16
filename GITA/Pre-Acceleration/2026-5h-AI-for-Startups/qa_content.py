"""Comprehensive content completeness audit.
For each slide, check whether all source blocks have rendered text.
"""
import json
from pptx import Presentation
from pptx.util import Emu


def emu_to_in(v):
    return v / 914400 if v else 0


def audit(path, json_path, label):
    print(f"\n{'='*70}\n  {label}\n{'='*70}")
    prs = Presentation(path)
    src = json.load(open(json_path, encoding='utf-8'))

    issues = []
    for i, (pslide, jslide) in enumerate(zip(prs.slides, src), 1):
        # Collect all PPTX text
        pptx_text = ""
        for sh in pslide.shapes:
            if sh.has_text_frame:
                pptx_text += sh.text_frame.text + " "
            elif sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        pptx_text += cell.text_frame.text + " "

        # Collect all source text
        def collect_source_text(b):
            t = b.get("type")
            text = b.get("text", "")
            if t == "list":
                items = b.get("items", [])
                text = " ".join(items)
            elif t == "table":
                rows = b.get("rows", [])
                text = " ".join(" ".join(r) for r in rows)
            elif t == "code":
                text = b.get("text", "")
            elif t == "box":
                sub_texts = []
                if b.get("text"):
                    sub_texts.append(b["text"])
                for sb in b.get("blocks", []) or []:
                    sub_texts.append(collect_source_text(sb))
                text = " ".join(filter(None, sub_texts))
            return text

        source_text = ""
        for b in jslide.get("blocks", []):
            source_text += collect_source_text(b) + " "

        # Calculate missing characters
        src_words = set(source_text.lower().split())
        pptx_words = set(pptx_text.lower().split())
        # Filter out common noise
        noise = {"the", "a", "an", "in", "on", "of", "to", "and", "or", "is", "are", "with", "/", "·"}
        meaningful_src = src_words - noise
        meaningful_pptx = pptx_words - noise

        if meaningful_src:
            coverage = len(meaningful_src & meaningful_pptx) / len(meaningful_src)
            if coverage < 0.85:
                missing = meaningful_src - meaningful_pptx
                # Filter very short/duplicate words
                missing = [w for w in missing if len(w) > 2][:10]
                issues.append((i, jslide.get("title", "")[:40], coverage, missing))

    if not issues:
        print(f"  ✓ All {len(src)} slides have ≥85% source content rendered.")
    else:
        print(f"  ⚠ {len(issues)} slides with potential content gaps:")
        for idx, title, cov, missing in issues:
            print(f"    Slide {idx:2d} ({cov*100:5.1f}%): {title!r}")
            if missing:
                print(f"      missing words: {missing}")


audit('day1-slides.pptx', 'day1-slides.json', 'Day 1 — AI Foundation & Validation')
audit('day2-slides.pptx', 'day2-slides.json', 'Day 2 — Product & Launch')